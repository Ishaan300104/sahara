"""Build Sahara final PPT by filling the course template with project content."""
from pptx import Presentation
from copy import deepcopy
from lxml import etree

SRC = "EAI Course Project Presentation Template.pptx"
OUT = "Sahara_Final_Prototype_Presentation.pptx"

# (slide_index_1based, old_text, new_text). old_text must match paragraph's
# full plain text. If old_text not found, we fall through silently — the
# matcher below prints warnings so we can fix mismatches.
REPLACEMENTS = [
    # Slide 1 — Title
    (1, "[Your Application Name]", "Sahara — सहारा"),
    (1, "Ethical Issues with AI — [Course Code]",
        "Ethical Issues with AI — MSL7500"),
    (1, "[Name 1, Name 2, Name 3, Name 4]",
        "Ishaan Pandey (B22CI017) · Om Singh (B22ME044) · Deekshant Singh Rajawat (B22CI010)"),
    (1, "Date: [Submission Date]  |  Deadline: 15 April 2026",
        "Date: 21 April 2026  |  IIT Jodhpur · School of Management and Entrepreneurship"),

    # Slide 3 — Problem Statement
    (3, "[Describe the real-world problem your application addresses. Be specific about who is affected and the scale of the problem.]",
        "Indian college students — especially at high-pressure technical institutes like IIT Jodhpur — face rising stress, anxiety, exam pressure and loneliness. Stigma, limited counselor availability, and long wait times stop many from seeking help early. The problem is large, time-sensitive, and often first surfaces at odd hours when no counselor is reachable."),
    (3, "[Explain why AI is the right approach for this problem, as opposed to conventional software.]",
        "AI is always-available, non-judgmental, and low-friction — a first-contact listener a student can reach at 2 a.m. Unlike a static FAQ app, an LLM can hold an empathetic, contextual conversation and adapt to how the student is actually feeling."),
    (3, "[Define your primary and secondary user groups]",
        "Primary: undergraduate students at IIT Jodhpur. Secondary: the institute Student Wellness Cell (for anonymised aggregate trends) and peer student populations at similar technical institutes."),

    # Slide 4 — Application Overview
    (4, "[What the app does end-to-end]",
        "A chat-first peer-listener web app. The student opens the site, consents, types how they feel, and Sahara responds empathetically. Mood is optionally logged. On crisis signals, the AI is bypassed and real helpline contacts are shown instantly."),
    (4, "[Frontend, backend, AI model/API, database]",
        "Flask + Jinja backend; vanilla HTML/CSS/JS + Chart.js frontend. Groq API serving Llama-3.3-70B-Versatile. Session-only state (Flask session) — no persistent database. Hosted on Render."),
    (4, "[What data the app uses and where it comes from]",
        "Only in-session user text + an optional user-selected mood tag. No names, no emails, no identifiers. No external datasets. Nothing persists after the session or the 🗑️ reset."),
    (4, "[Specifically how AI is used: classification, generation, recommendation, etc.]",
        "Generation: empathetic natural-language replies via LLM. Rule-based classification: a 30-phrase safety filter runs before the LLM to catch crisis language. The AI does not diagnose, prescribe, or store anything."),

    # Slide 5 — Architecture & Screenshots
    (5, "[Insert System Architecture Diagram]",
        "5-Layer Architecture: (1) Interface  →  (2) AI Brain (Llama 3.3 via Groq)  →  (3) Safety Filter (keyword match, bypasses LLM on crisis)  →  (4) Mood Tracker (opt-in, client-side)  →  (5) Data Storage (Flask session only, no DB)"),
    (5, "[Insert Key Application Screenshots]",
        "Screens demonstrated in live demo: Consent modal · Chat with AI-label + ℹ️ Why · 🚩 Report button · Mood chart · Crisis banner with iCall / NIMHANS / IITJ counselor · Resources tab · Breathing exercise"),

    # Slide 6 — Ethical Considerations in Design
    (6, "[How you minimized data collection, implemented consent, secured storage]",
        "Session-only state; no login; no personal identifiers collected. Explicit checkbox consent before the first chat. One-click 🗑️ reset wipes everything. No third-party analytics."),
    (6, "[Steps taken to identify and reduce bias in training data, model outputs]",
        "System prompt anchors a warm, non-judgmental peer tone and forbids diagnosis or medication talk. Safety filter is a transparent phrase list — auditable. Known limitation acknowledged: English-only and Western framing of mental-health concepts."),
    (6, "[How users can understand AI decisions; explainability features built in]",
        "Persistent \"⚠️ AI — Not a therapist\" badge. Every bot reply now carries a \"🌊 Sahara · AI peer listener\" label and an ℹ️ \"Why this response?\" popup explaining the model, context, and the Sahara system prompt used."),
    (6, "[Opt-in/opt-out mechanisms, user control over AI actions]",
        "Consent is required before the first message. Mood tracking is fully optional. User can end / clear the session anytime via 🗑️. A 🚩 button lets the user reject any reply they find unhelpful or off."),
    (6, "[Guardrails, content filtering, edge case handling]",
        "Crisis filter runs BEFORE the LLM — on a hit, the LLM is bypassed and iCall (9152987821), NIMHANS, Vandrevala (24×7) and the IITJ counselor are shown immediately with a warm message, not clinical jargon."),

    # Slide 7 — Feedback process steps
    (7, "[Who, how many, selection criteria]",
        "10 student respondents — mix of IIT Jodhpur UG peers and team members. Selected for a spread across familiarity with AI (Basic / Intermediate / Advanced) and usage duration (10–20 min)."),
    (7, "[Format: demo + questionnaire, interview, etc.]",
        "Remote hands-on prototype walkthrough via the public web link, followed by an individually-filled questionnaire. Everyone tried the consent flow, a normal chat, the mood tracker and the crisis-banner scenario."),
    (7, "[Questionnaire covering usability, AI quality, ethics]",
        "30-question questionnaire covering: usability (Section B), AI output quality (Section C), ethics — privacy, bias, transparency, consent, safety, societal impact (Section D) and overall assessment (Section E)."),
    (7, "[Quantitative ratings + qualitative themes]",
        "Likert-scale ratings aggregated per section. Free-text answers coded into themes: transparency, accountability, safety escalation, language bias, data control, UI polish."),

    # Slide 8 — Quantitative (replace bullet placeholders with real numbers)
    (8, "[Insert charts/graphs showing aggregated rating scale responses from the feedback questionnaire. Include at minimum:]",
        "Aggregated Likert ratings across N = 10 respondents (scale 1–5):"),
    (8, "•  Average usability rating across all respondents",
        "•  Purpose clarity (Q1): 4.4/5  ·  UI rating (Q2): 4.1/5  ·  Likely-to-use (Q5): 4.1/5"),
    (8, "•  AI accuracy/relevance rating distribution",
        "•  AI accuracy/relevance (Q6): 3.7/5  ·  5 × Accurate, 4 × Mixed, 1 × Very Accurate"),
    (8, "•  Ethical responsibility rating distribution",
        "•  Ethical responsibility (Q28): 3.7/5  ·  7 × Responsible, 3 × Neutral, 0 × Irresponsible"),
    (8, "•  Data comfort level distribution",
        "•  Data comfort (Q11): 3.8/5  ·  Control over own data (Q22): 2.7/5 ← 2nd lowest"),
    (8, "•  Fairness confidence rating distribution",
        "•  Fairness confidence (Q16): 3.5/5  ·  Explainability (Q8): 2.3/5 ← LOWEST"),
    (8, "[Insert your charts here — bar charts, radar charts, or stacked bar charts work well]",
        "Key takeaway: users trust the concept and interface, but explainability (2.3/5) and data control (2.7/5) are the two clear weak spots driving our iterations."),

    # Slide 9 — Qualitative Themes
    (9, "[e.g., 6/12 respondents found navigation confusing on the results page]",
        "2/10 (Tejas, Avichal) noted slow response loading. 1/10 (Ishaan) found mood history hard to locate at first. 1/10 (Tejas) felt the UI was a bit plain and would hurt trust."),
    (9, "[e.g., 4/12 respondents reported irrelevant recommendations in specific scenarios]",
        "2/10 (Deekshant, Maulik) got a reply that felt tone-deaf / dismissive of a serious concern. 2/10 (Om, Amey) said replies were occasionally generic or templated. 1/10 (Maulik) flagged minor hallucination (advice not matching emotional context)."),
    (9, "[e.g., 8/12 respondents were concerned about data retention; 5/12 noted lack of AI explainability]",
        "8/10 said there was no way to report a bad AI response (Q18). 7/10 could not distinguish AI-generated from scripted replies (Q17). 5/10 said English-only / Western framing creates cultural-fit gaps. 4/10 wanted privacy disclaimers repeated mid-chat, not only at start."),
    (9, "[e.g., 7/12 respondents wanted ability to export their data]",
        "2/10 (Ankit, Deekshant) asked for Hindi / regional language support. 2/10 (Ishaan, Tejas, Avichal) wanted multi-session mood trends / session history. 1/10 (Amey) asked for an in-app rating of each response."),
    (9, "[e.g., 10/12 found the core concept valuable; 9/12 liked the interface design]",
        "10/10 found the core concept valuable. 8/10 rated UI Good or Excellent. 8/10 rated the app ethically Responsible. 9/10 said they would likely use it if fully developed. Multiple respondents singled out the empathetic, non-judgmental tone as a real strength."),

    # Slide 10 — Ethical Issues Identified from Feedback
    (10, "Issue 1: [e.g., Privacy / Excessive Data Collection]",
        "Issue 1: Lack of Explainability of AI Decisions"),
    (10, "Evidence: [X/N respondents raised this]. Specific feedback: [quote or paraphrase].",
        "Evidence: Q8 explainability averaged 2.3/5 (lowest of all ratings). 8/10 rated it Poor. 7/10 (Q17) could not tell AI-generated from scripted replies. No in-app \"why this response\" path existed."),
    (10, "Issue 2: [e.g., Bias in AI Outputs]",
        "Issue 2: No Error-Reporting / Accountability Mechanism"),
    (10, "Evidence: [X/N respondents observed this]. Specific example: [describe the biased output].",
        "Evidence: 8/10 said there was no way to flag a bad response (Q18). 2/10 (Deekshant, Maulik) described a specific tone-deaf reply they had no way to report. Amey explicitly asked for a response-rating feature."),
    (10, "Issue 3: [e.g., Lack of Transparency / Explainability]",
        "Issue 3: Insufficient Safety Escalation for Vulnerable Users"),
    (10, "Evidence: [X/N respondents could not understand how AI arrived at decisions].",
        "Evidence: 7/10 chose crisis-escalation as the single most important ethical improvement (Q29). 10/10 named users in acute distress / severe mental illness as most at risk (Q25). \"False reassurance\" was the most-cited harm scenario."),
    (10, "Issue 4: [e.g., Insufficient Consent Mechanisms]",
        "Issue 4: Language and Cultural Bias"),
    (10, "Evidence: [X/N respondents did not recall being asked for consent].",
        "Evidence: 5/10 (Om, Deekshant, Maulik, Tejas, Ankit) noted English-only or Western-centric framing creates equity gaps. Fairness confidence (Q16) averaged only 3.5/5. Ankit explicitly asked for Hindi support."),
    (10, "Issue 5: [e.g., Potential for Harm to Vulnerable Users]",
        "Issue 5: Weak Consent Depth & Data Control"),
    (10, "Evidence: [X/N respondents identified specific harm scenarios].",
        "Evidence: Control over own data (Q22) averaged only 2.7/5 (2nd lowest rating). 4/10 said the start-of-session disclaimer is not enough and wanted a fuller, repeated policy. Maulik called the one-time pop-up \"not truly informed consent.\""),
    (10, "Add more slides if you have more than 5 ethical issues. Each issue must have respondent evidence.",
        "These 5 issues map directly to the rows on the next slide — each has a concrete change shipped in the final prototype."),

    # Slide 12 — Before & After
    (12, "[Insert screenshot or description of how the application looked/behaved before ethical improvements]",
        "BEFORE: One start-of-session disclaimer. Bot replies had no AI label. No way to flag an unhelpful reply. \"Talk to Counselor\" was a faint outline button in the nav. Consent modal text was short. No inline explainability."),
    (12, "[Insert screenshot or description of the improved version after addressing ethical issues from feedback]",
        "AFTER: Every bot reply shows \"🌊 Sahara · AI peer listener · not a therapist\" with 🚩 Report and ℹ️ Why this response? buttons. Disclaimer reminder toast auto-appears every 5 messages. Talk to Counselor is now a solid, shadowed primary-colour button. Consent modal explicitly states session-only + no identifiers + 🗑️ reset."),

    # Slide 13 — Demo
    (13, "[Live Demo or Recorded Walkthrough of the Final Working Prototype]",
        "Live walkthrough of the final Sahara prototype — public web link shown during demo. Tested on desktop and mobile."),

    # Slide 14 — Reflection
    (14, "[What were the hardest technical problems to solve? How did you resolve them?]",
        "Calibrating the safety filter to catch implicit distress phrasing (e.g. \"can't go on\") without over-triggering on casual stress. Designing session-only state so nothing leaks across users without adding a database. Groq rate-limit handling and keeping context short for cost + privacy."),
    (14, "[What ethical issues surprised you? What did you learn about responsible AI development?]",
        "Users care far more about transparency and accountability than about fancy features. A single consent pop-up is not enough — reinforcement during the conversation matters. Empathetic tone alone does not justify the risk of false reassurance in a mental-health context."),
    (14, "[What worked well in your feedback collection? What would you change?]",
        "The structured 30-question questionnaire made answers comparable and surfaced the two weak spots (explainability 2.3/5, data control 2.7/5) cleanly. Next time: pre-brief respondents on what a good crisis-response should look like, so Q29 answers go deeper than \"make it more prominent.\""),
    (14, "[What ethical issues remain unresolved or partially addressed? Be honest about limitations.]",
        "English-only language limits reach across non-urban students. No real-time human-counselor handoff yet. LLM hallucination and tone-deaf-response risk remains. Environmental cost of LLM inference is not addressed in-app."),
    (14, "[What further ethical improvements would you make with additional time and resources?]",
        "Hindi + regional language support (2 respondents asked). Real counselor tele-handoff from the crisis banner. Source citations for coping-strategy advice. Journaling + weekly mood-summary features requested by users. A full DPDP-aligned privacy policy page."),

    # Slide 15 — Thank You
    (15, "[Name 1, Name 2, Name 3, Name 4]",
        "Ishaan Pandey · Om Singh · Deekshant Singh Rajawat"),
    (15, "[Application Name]",
        "Sahara — सहारा · AI-powered mental-health peer-listener chatbot for college students"),
    (15, "[Email or link to prototype]",
        "Prototype: public Render link shown in demo · Course: MSL7500 — Ethical Issues with AI · IIT Jodhpur"),
]

# Table replacements: (slide_idx_1based, table_shape_name, row, col, new_text)
TABLE_REPLACEMENTS = [
    # Slide 7 — demographics table
    (7, "Table 0", 1, 1, "10"),
    (7, "Table 0", 2, 1,
        "IIT Jodhpur UG students. Familiarity with AI: 3 Basic / 6 Intermediate / 1 Advanced."),
    (7, "Table 0", 3, 1, "18 April 2026 – 21 April 2026"),
    (7, "Table 0", 4, 1,
        "Remote hands-on web demo (10–20 min) followed by a structured 30-question questionnaire across usability, AI quality and ethics."),

    # Slide 11 — iterations table (header + 4 issue rows)
    (11, "Table 0", 1, 0, "Lack of AI transparency / explainability"),
    (11, "Table 0", 1, 1, "Q8 averaged 2.3/5; 8/10 rated Poor. 7/10 could not tell AI vs scripted (Q17)."),
    (11, "Table 0", 1, 2, "Every bot reply now shows a \"🌊 Sahara · AI peer listener · not a therapist\" label and an ℹ️ \"Why this response?\" popup explaining model + context used."),
    (11, "Table 0", 1, 3, "Visible on every bot reply in final prototype; tested by team."),

    (11, "Table 0", 2, 0, "No way to report a bad AI response"),
    (11, "Table 0", 2, 1, "8/10 said no mechanism exists (Q18). 2/10 had a specific tone-deaf reply they could not flag."),
    (11, "Table 0", 2, 2, "Added 🚩 flag button on every bot message, POST /report endpoint logs flagged replies for review."),
    (11, "Table 0", 2, 3, "Live-testable in demo — 🚩 turns to ✓ on click."),

    (11, "Table 0", 3, 0, "Weak disclaimer + crisis escalation"),
    (11, "Table 0", 3, 1, "7/10 said crisis escalation is the top improvement (Q29); 4/10 wanted periodic disclaimers (Q19)."),
    (11, "Table 0", 3, 2, "Promoted \"📞 Talk to Counselor\" button to solid primary colour; auto-show AI-disclaimer reminder toast every 5 messages mid-chat."),
    (11, "Table 0", 3, 3, "Reminder now fires inside the conversation, not only at session start."),

    (11, "Table 0", 4, 0, "Shallow consent + data control"),
    (11, "Table 0", 4, 1, "Q22 control averaged 2.7/5. 4/10 wanted fuller policy. Maulik: \"one-time pop-up, not truly informed.\""),
    (11, "Table 0", 4, 2, "Tightened consent-modal text — explicitly states session-only data, no identifiers, and the 🗑️ reset button; reset remains one-click."),
    (11, "Table 0", 4, 3, "Consent screen re-reviewed with team pre-demo."),
]


def set_paragraph_text(para, new_text):
    """Replace paragraph text, preserving the formatting of its first run.

    Strategy: copy the first run, clear the paragraph of all runs, then add
    back a single run with the new text and the copied formatting.
    """
    runs = para.runs
    if not runs:
        # add a fresh run if para has none
        run = para.add_run()
        run.text = new_text
        return

    # preserve rPr of first run
    first_run = runs[0]
    r_elem = first_run._r
    # remove all existing <a:r> elements from the paragraph
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for r in para._pPr.getparent().findall("a:r", nsmap) if para._pPr is not None else para._p.findall("a:r", nsmap):
        para._p.remove(r)
    # also remove endParaRPr/br that aren't runs — keep them, just append new run before endParaRPr
    # Re-add the preserved run, with new text
    para._p.append(r_elem)
    first_run.text = new_text


def replace_paragraph_preserving_format(para, new_text):
    """Simpler, safer: keep first run's rPr, rewrite its text, delete the rest."""
    runs = para.runs
    if not runs:
        run = para.add_run()
        run.text = new_text
        return
    first_run = runs[0]
    first_run.text = new_text
    # remove any additional runs after the first
    r_elements = para._p.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
    )
    for r in r_elements[1:]:
        para._p.remove(r)


def iter_paragraphs(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                yield shape, para
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        yield shape, para


def main():
    prs = Presentation(SRC)

    # group replacements by slide
    by_slide = {}
    for sid, old, new in REPLACEMENTS:
        by_slide.setdefault(sid, []).append((old, new))

    for sid in sorted(by_slide):
        slide = prs.slides[sid - 1]
        pairs = by_slide[sid]
        matched = {old: False for old, _ in pairs}
        for shape, para in iter_paragraphs(slide):
            ptext = "".join(r.text for r in para.runs)
            for old, new in pairs:
                if ptext.strip() == old.strip():
                    replace_paragraph_preserving_format(para, new)
                    matched[old] = True
                    break
        for old, ok in matched.items():
            if not ok:
                print(f"WARN: slide {sid}: no paragraph matched: {old[:80]!r}")

    # tables
    for sid, tname, r, c, new_text in TABLE_REPLACEMENTS:
        slide = prs.slides[sid - 1]
        for shape in slide.shapes:
            if shape.has_table and shape.name == tname:
                cell = shape.table.rows[r].cells[c]
                # clear cell's paragraphs but keep first, replace text
                tf = cell.text_frame
                first_para = tf.paragraphs[0]
                replace_paragraph_preserving_format(first_para, new_text)
                # drop any subsequent paragraphs
                p_elems = tf._txBody.findall(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                )
                for p in p_elems[1:]:
                    tf._txBody.remove(p)
                break

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
